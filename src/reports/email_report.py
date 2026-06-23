import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import pandas as pd
import psycopg2
from pptx import Presentation
from pptx.util import Inches
from datetime import datetime
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email.mime.text import MIMEText
from email import encoders
import os

# ==========================================
# CONFIGURATION (!!! EDIT THESE !!!)
# ==========================================
DB_ENDPOINT = "lending-analytics-db.cotq60e68s8c.us-east-1.rds.amazonaws.com"  # e.g., lending-analytics-db.xxxxxx.ca-central-1.rds.amazonaws.com
DB_NAME = "lendingdb"
DB_USER = "lending_admin"
DB_PASSWORD = "LendingSecurePass123!"

# EMAIL CONFIG (Use a Gmail "App Password" - see guide below)
SMTP_SERVER = "smtp.gmail.com"
SMTP_PORT = 587
SENDER_EMAIL = "ddeepachandrasekar30@gmail.com"   
SENDER_PASSWORD = "maro qkod gggx hdbj"   
RECEIVER_EMAIL = "deepachandrasekar30@gmail.com" 

# ==========================================
# FUNCTION: Generate the PPT Report
# ==========================================
def generate_ppt():
    print("📊 Connecting to AWS RDS...")
    conn = psycopg2.connect(
        host=DB_ENDPOINT,
        database=DB_NAME,
        user=DB_USER,
        password=DB_PASSWORD,
        port=5432
    )
    
    # Load data
    query = """
        SELECT 
            p.loan_id,
            p.remaining_balance,
            p.payment_status,
            c.province,
            pr.product_type
        FROM loan_portfolio p
        JOIN clients c ON p.client_id = c.client_id
        JOIN loan_products pr ON p.product_id = pr.product_id
    """
    df = pd.read_sql(query, conn)
    conn.close()
    print(f"✅ Loaded {len(df)} records from cloud DB.")

    # Calculate metrics
    total_loans = len(df)
    total_value = df['remaining_balance'].sum()
    delinquent = df[df['payment_status'].isin(['90+ Days Past Due', 'Default'])]
    del_rate = (len(delinquent) / total_loans) * 100 if total_loans > 0 else 0

    # Build Chart: Delinquency by Province
    prov_del = df[df['payment_status'].isin(['90+ Days Past Due', 'Default'])].groupby('province').size()
    
    fig, ax = plt.subplots(figsize=(10, 6))
    if not prov_del.empty:
        prov_del.plot(kind='bar', ax=ax, color='#d4af37')
        ax.set_title('Delinquent Loans by Province', fontsize=16)
        ax.set_ylabel('Count')
        ax.grid(axis='y', linestyle='--', alpha=0.7)
    else:
        ax.text(0.5, 0.5, 'No Delinquent Loans', ha='center', va='center')
    
    chart_path = 'data/delinquency_chart.png'
    plt.savefig(chart_path, bbox_inches='tight')
    plt.close()

    # Create PowerPoint
    print("📁 Generating PowerPoint...")
    prs = Presentation()
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Lending Portfolio Health Report"
    slide.placeholders[1].text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}"

    # Slide 2: KPIs
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Summary"
    text_frame = slide.placeholders[1].text_frame
    text_frame.text = f"Total Active Loans: {total_loans:,}\n"
    text_frame.text += f"Total Portfolio Value: ${total_value:,.2f}\n"
    text_frame.text += f"Delinquency Rate: {del_rate:.2f}%\n"
    text_frame.text += f"At-Risk Loans: {len(delinquent)}"

    # Slide 3: Chart
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Regional Delinquency Breakdown"
    left = Inches(1)
    top = Inches(1.5)
    slide.shapes.add_picture(chart_path, left, top, width=Inches(8))

    output_ppt = 'data/Lending_Portfolio_Report.pptx'
    prs.save(output_ppt)
    print(f"✅ PPT saved to {output_ppt}")
    return output_ppt

# ==========================================
# FUNCTION: Send Email
# ==========================================
def send_email(attachment_path):
    print("📧 Preparing to send email...")
    
    msg = MIMEMultipart()
    msg['From'] = SENDER_EMAIL
    msg['To'] = RECEIVER_EMAIL
    msg['Subject'] = f"Daily Lending Report - {datetime.now().strftime('%Y-%m-%d')}"

    body = "Attached is the daily lending portfolio health report."
    msg.attach(MIMEText(body, 'plain'))

    # Attach file
    try:
        with open(attachment_path, "rb") as attachment:
            part = MIMEBase('application', 'octet-stream')
            part.set_payload(attachment.read())
            encoders.encode_base64(part)
            part.add_header(
                'Content-Disposition',
                f"attachment; filename= {os.path.basename(attachment_path)}",
            )
            msg.attach(part)
    except Exception as e:
        print(f"❌ Failed to attach file: {e}")
        return

    # Send Email
    try:
        server = smtplib.SMTP(SMTP_SERVER, SMTP_PORT)
        server.starttls()
        server.login(SENDER_EMAIL, SENDER_PASSWORD)
        server.sendmail(SENDER_EMAIL, RECEIVER_EMAIL, msg.as_string())
        server.quit()
        print("✅ Email sent successfully!")
    except Exception as e:
        print(f"❌ Email failed: {e}")

# ==========================================
# MAIN EXECUTION
# ==========================================
if __name__ == "__main__":
    # 1. Generate the PPT
    ppt_file = generate_ppt()
    
    # 2. Send the email
    send_email(ppt_file)